"""
Build SF Tech Debt Assessor executive presentation.
Run: python3 docs/build_presentation.py
Output: docs/SF_Tech_Debt_Assessor_Executive_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
SF_BLUE       = RGBColor(0x00, 0xA1, 0xE0)   # Salesforce blue
SF_DARK       = RGBColor(0x16, 0x25, 0x5E)   # Salesforce dark navy
SF_LIGHT_BLUE = RGBColor(0xD0, 0xEE, 0xF8)   # light blue background accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY      = RGBColor(0x54, 0x54, 0x54)
LIGHT_GREY    = RGBColor(0xF3, 0xF4, 0xF6)
CRITICAL_RED  = RGBColor(0xC2, 0x39, 0x34)
HIGH_ORANGE   = RGBColor(0xE8, 0x77, 0x22)
MEDIUM_YELLOW = RGBColor(0xF5, 0xBE, 0x26)
LOW_GREEN     = RGBColor(0x2E, 0x9E, 0x6B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def header_bar(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill=SF_DARK)
    add_text(slide, title, 0.4, 0.1, 10, 0.6,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.65, 10, 0.4,
                 size=14, color=SF_LIGHT_BLUE)
    # Blue accent line under bar
    add_rect(slide, 0, 1.1, 13.33, 0.06, fill=SF_BLUE)


def footer(slide, page_num):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=SF_DARK)
    add_text(slide, "SF Tech Debt Assessor  |  Confidential  |  Steven Bilgram, Success Architect",
             0.3, 7.12, 10, 0.28, size=9, color=SF_LIGHT_BLUE)
    add_text(slide, str(page_num), 12.8, 7.12, 0.4, 0.28,
             size=9, color=SF_LIGHT_BLUE, align=PP_ALIGN.RIGHT)


def bullet_box(slide, items, l, t, w, h,
               title=None, title_color=SF_DARK,
               bullet="•", size=13, bg=None, border=None):
    if bg:
        add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=0.5)
    txb = slide.shapes.add_textbox(Inches(l + 0.12), Inches(t + 0.1),
                                   Inches(w - 0.24), Inches(h - 0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = title
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = title_color
    for i, item in enumerate(items):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = DARK_TEXT


def severity_pill(slide, label, color, l, t):
    add_rect(slide, l, t, 1.5, 0.32, fill=color)
    add_text(slide, label, l, t, 1.5, 0.32,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full background gradient simulation — dark navy
add_rect(slide, 0, 0, 13.33, 7.5, fill=SF_DARK)
# Blue diagonal accent
add_rect(slide, 8.5, 0, 4.83, 7.5, fill=SF_BLUE)
add_rect(slide, 7.8, 0, 1.2, 7.5, fill=SF_DARK)  # overlap to soften edge

# Logo area placeholder
add_rect(slide, 0.5, 0.5, 2.2, 0.7, fill=SF_BLUE)
add_text(slide, "Salesforce", 0.55, 0.52, 2.1, 0.6,
         size=16, bold=True, color=WHITE)

add_text(slide, "Salesforce", 0.4, 1.8, 8, 0.7,
         size=20, color=SF_LIGHT_BLUE)
add_text(slide, "Tech Debt Assessor", 0.4, 2.35, 8.5, 1.1,
         size=52, bold=True, color=WHITE)
add_text(slide, "A comprehensive, read-only org health assessment across\n341 checks in 23 categories — surfacing technical debt,\nsecurity gaps, and actionable remediation priorities.",
         0.4, 3.6, 7.5, 1.4, size=16, color=SF_LIGHT_BLUE)

add_text(slide, "Steven Bilgram  |  Success Architect", 0.4, 5.5, 6, 0.4,
         size=14, color=WHITE)
add_text(slide, "June 2026", 0.4, 5.9, 4, 0.4,
         size=13, color=SF_LIGHT_BLUE, italic=True)

# Right side stats
for i, (num, lbl) in enumerate([("341", "Total Checks"), ("23", "Categories"),
                                  ("Read-Only", "Zero Org Impact"), ("< 60s", "Full Assessment")]):
    y = 1.5 + i * 1.3
    add_rect(slide, 9.3, y, 3.3, 1.1, fill=WHITE)
    add_text(slide, num, 9.3, y + 0.05, 3.3, 0.6,
             size=32, bold=True, color=SF_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, lbl, 9.3, y + 0.62, 3.3, 0.4,
             size=12, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "The Problem", "Salesforce orgs accumulate technical debt silently — until it costs you")
footer(slide, 2)

pain_points = [
    ("Security gaps go undetected", "OWD misconfiguration, guest access, missing MFA, stale OAuth tokens — none visible without a deep audit."),
    ("Platform deprecations break releases", "SOAP login(), retired API versions, Session IDs in Outbound Messages, PushTopics — all with hard enforcement dates."),
    ("Automation complexity compounds", "Workflow Rules, Process Builders, and Flows overlap with no single view of routing, coverage, or risk."),
    ("Test coverage is misleading", "Managed package classes inflate coverage numbers. Abstract classes with zero executable lines show as untested."),
    ("Consultant assessments take days", "Manual review across 23 areas — org limits, sharing rules, flows, Knowledge, entitlements — takes 2–3 days per engagement."),
]

for i, (title, desc) in enumerate(pain_points):
    col = i % 3
    row = i // 3
    l = 0.3 + col * 4.35
    t = 1.35 + row * 2.7
    add_rect(slide, l, t, 4.1, 2.4, fill=LIGHT_GREY, line=SF_BLUE, line_w=0.5)
    add_rect(slide, l, t, 4.1, 0.38, fill=SF_DARK)
    add_text(slide, title, l + 0.1, t + 0.04, 3.9, 0.32,
             size=12, bold=True, color=WHITE)
    add_text(slide, desc, l + 0.12, t + 0.48, 3.86, 1.82,
             size=12, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What It Is
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "What the SF Tech Debt Assessor Does", "Connect any org. Get a full health report in under 60 seconds.")
footer(slide, 3)

add_text(slide, "One OAuth connection. No data written to Salesforce. All metadata is ephemeral.",
         0.4, 1.3, 12.5, 0.45, size=15, color=MID_GREY, italic=True)

steps = [
    ("1", "Connect", "Enter org URL, Client ID, Client Secret.\nOAuth login — same as any Salesforce app."),
    ("2", "Scan", "341 checks run automatically across\n23 categories in parallel. ~45 seconds."),
    ("3", "Review", "Colour-coded category grid with scores.\nExpand any finding to see affected records."),
    ("4", "Export", "PDF report, Excel workbook, CSV flat file,\nor Remediation Roadmap (print to PDF)."),
]

arrow_pts = [3.55, 6.45, 9.35]
for x in arrow_pts:
    add_text(slide, "→", x, 2.6, 0.5, 0.7, size=28, color=SF_BLUE, align=PP_ALIGN.CENTER)

for i, (num, title, desc) in enumerate(steps):
    l = 0.3 + i * 3.2
    add_rect(slide, l, 2.0, 2.9, 3.3, fill=SF_DARK)
    add_text(slide, num, l + 0.1, 2.1, 0.6, 0.6,
             size=32, bold=True, color=SF_BLUE)
    add_text(slide, title, l + 0.1, 2.65, 2.7, 0.45,
             size=17, bold=True, color=WHITE)
    add_text(slide, desc, l + 0.1, 3.15, 2.7, 1.9,
             size=12, color=SF_LIGHT_BLUE)

add_rect(slide, 0.3, 5.55, 12.7, 1.2, fill=SF_LIGHT_BLUE)
props = [
    "Read-only — zero writes to Salesforce",
    "No database — results live only in the browser tab",
    "OAuth tokens expire after 1 hour",
    "Hosted at sf-tech-debt-assessor.onrender.com",
    "Docker image available for on-premises / VPN use",
]
txb = slide.shapes.add_textbox(Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.0))
txb.word_wrap = False
tf = txb.text_frame
p = tf.paragraphs[0]
for j, prop in enumerate(props):
    if j > 0:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = ("  |  " if j > 0 else "") + "✓  " + prop
    run.font.size = Pt(12)
    run.font.color.rgb = SF_DARK
    run.font.bold = (j == 0)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 23 Categories Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "23 Assessment Categories — 341 Checks", "Every category contributes to an overall org health score")
footer(slide, 4)

categories = [
    ("Configuration", "13"),
    ("Code Quality", "46"),
    ("Data Model", "4"),
    ("Service Cloud", "69"),
    ("Sharing & Security", "24"),
    ("Integrations", "9"),
    ("Test Coverage", "7"),
    ("Org Limits", "5"),
    ("Duplicate Rules", "4"),
    ("Reports & Dashboards", "3"),
    ("Email Templates", "3"),
    ("Platform Events & CDC", "3"),
    ("Managed Packages", "3"),
    ("Custom Metadata", "3"),
    ("Record Types & Layouts", "4"),
    ("Einstein & AI", "9"),
    ("Experience Cloud", "15"),
    ("Connected App Security", "12"),
    ("LWC & Components", "39"),
    ("OmniStudio", "26"),
    ("Performance", "22"),
    ("Notes & Attachments", "12"),
    ("Flow Quality", "6"),
]

cols = 6
col_w = 13.1 / cols
row_h = 0.72
for i, (name, count) in enumerate(categories):
    col = i % cols
    row = i // cols
    l = 0.12 + col * col_w
    t = 1.3 + row * (row_h + 0.08)
    bg = SF_DARK if i < 3 else (SF_BLUE if i < 8 else LIGHT_GREY)
    txt_color = WHITE if bg in (SF_DARK, SF_BLUE) else DARK_TEXT
    cnt_color = SF_BLUE if bg == SF_DARK else (WHITE if bg == SF_BLUE else SF_DARK)
    add_rect(slide, l, t, col_w - 0.1, row_h, fill=bg)
    add_text(slide, count, l + 0.08, t + 0.02, 0.55, 0.38,
             size=20, bold=True, color=cnt_color)
    add_text(slide, name, l + 0.08, t + 0.36, col_w - 0.2, 0.34,
             size=9, color=txt_color)

add_text(slide, "Dark navy = most impactful categories  •  Blue = high check density  •  All 23 categories scored 0–100",
         0.3, 7.0, 12.7, 0.35, size=10, color=MID_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Service Cloud Deep Dive
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Service Cloud — 69 Checks", "The most comprehensive Service Cloud org audit available")
footer(slide, 5)

areas = [
    ("Case Configuration", ["Excessive queues (>25)", "Inactive record types", "Legacy assignment/escalation rules", "Unverified OWAs ⚠ Spring '26"]),
    ("Omni-Channel", ["No service channels configured", "Tab-based capacity (legacy)", "No push timeout on routing", "No presence configurations"]),
    ("Knowledge", ["No published articles", "Stalled draft articles (180+ days)", "Stale published articles (12+ months)", "No data category groups"]),
    ("Entitlements", ["No business hours on SLA processes", "No milestone actions", "Open cases with no SLA start", "Expired entitlements still active"]),
    ("Email-to-Case", ["Routing addresses without TLS", "No default case owner", "Email threading gaps", "Unrestricted email service addresses"]),
    ("SLA & Case Health", ["Active milestone violations", "Stale escalated cases", "High zero-touch close rate", "Cases with no contact/account"]),
    ("Live Chat / MIAW", ["Non-Omni-Channel chat buttons", "Legacy Live Chat — MIAW not adopted", "Chat buttons → deleted queues", "Messaging channels without OPTOUT"]),
    ("Agent Efficiency", ["No Service Console app", "No active macros", "No Quick Texts", "No Einstein NBA strategies"]),
]

cols = 4
cw = 3.2
ch = 2.4
for i, (title, checks) in enumerate(areas):
    col = i % cols
    row = i // cols
    l = 0.12 + col * (cw + 0.08)
    t = 1.28 + row * (ch + 0.1)
    add_rect(slide, l, t, cw, ch, fill=LIGHT_GREY)
    add_rect(slide, l, t, cw, 0.35, fill=SF_BLUE)
    add_text(slide, title, l + 0.1, t + 0.04, cw - 0.15, 0.3,
             size=11, bold=True, color=WHITE)
    for j, chk in enumerate(checks):
        add_text(slide, f"• {chk}", l + 0.1, t + 0.42 + j * 0.44, cw - 0.2, 0.42,
                 size=10, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Security & Platform Risk
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Security & Platform Risk", "Deprecation deadlines and security gaps that require immediate attention")
footer(slide, 6)

# Left — deprecation timeline
add_rect(slide, 0.25, 1.25, 6.0, 5.8, fill=LIGHT_GREY)
add_text(slide, "Breaking Changes — Enforcement Timeline", 0.4, 1.32, 5.7, 0.4,
         size=13, bold=True, color=SF_DARK)

timeline = [
    ("Feb 2026", "Session IDs in Outbound Messages retired", CRITICAL_RED, True),
    ("Mar 2026", "CA certificates max lifespan 200 days", CRITICAL_RED, True),
    ("Spring '26", "SOAP login() disabled for new orgs", HIGH_ORANGE, True),
    ("Spring '26", "My Domain login URL enforced (production)", HIGH_ORANGE, True),
    ("Spring '26", "Unverified OWAs fail to send", HIGH_ORANGE, True),
    ("May 2026", "Phishing-resistant MFA for privileged users", CRITICAL_RED, True),
    ("Summer '26", "WCAG 2.2 force-applied to Experience Cloud", MEDIUM_YELLOW, False),
    ("Summer '26", "PushTopics (Streaming API) deprecated", HIGH_ORANGE, False),
    ("Summer '27", "SOAP login() hard retirement", CRITICAL_RED, False),
]

for i, (date, label, color, past) in enumerate(timeline):
    y = 1.85 + i * 0.54
    add_rect(slide, 0.35, y, 1.5, 0.38, fill=color)
    add_text(slide, date, 0.35, y, 1.5, 0.38,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, ("⚠ " if past else "→ ") + label,
             1.97, y + 0.04, 4.1, 0.34,
             size=10, color=DARK_TEXT, bold=past)

# Right — security checks
add_rect(slide, 6.55, 1.25, 6.55, 5.8, fill=SF_DARK)
add_text(slide, "Sharing & Security Checks (24 total)", 6.7, 1.32, 6.2, 0.4,
         size=13, bold=True, color=WHITE)

sec_checks = [
    ("Critical", CRITICAL_RED, ["OWD Public Read/Write (internal & external)", "Non-admin users with Modify All Data", "Users with no MFA enrollment"]),
    ("High", HIGH_ORANGE, ["Guest access to objects", "Cloned System Administrator profiles", "Transaction Security Policies not configured", "Privileged permission sets — broad access"]),
    ("Medium", MEDIUM_YELLOW, ["Users with password never expires", "Low-security sessions (> 20, MFA not enforced)", "OAuth tokens for deactivated users"]),
    ("Low", LOW_GREEN, ["Login IP restrictions not configured", "Users with excessive permission sets", "Async Sharing Recalculation not enabled"]),
]

y = 1.85
for sev, color, checks in sec_checks:
    add_rect(slide, 6.65, y, 1.1, 0.3, fill=color)
    add_text(slide, sev, 6.65, y, 1.1, 0.3,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for chk in checks:
        y += 0.36
        add_text(slide, f"• {chk}", 6.75, y, 6.2, 0.34,
                 size=10, color=SF_LIGHT_BLUE)
    y += 0.42


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Code & Flow Quality
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Code Quality, Test Coverage & Flow Quality", "46 + 7 + 6 checks — PMD, Graph Engine, and Flow Scanner rules")
footer(slide, 7)

# Code Quality column
add_rect(slide, 0.2, 1.28, 4.2, 5.85, fill=LIGHT_GREY)
add_rect(slide, 0.2, 1.28, 4.2, 0.42, fill=SF_DARK)
add_text(slide, "Code Quality — 46 Checks", 0.32, 1.3, 4.0, 0.38,
         size=13, bold=True, color=WHITE)

cq = [
    (CRITICAL_RED, "Critical (7)", ["DML/SOQL inside loops", "Dynamic SOQL injection risk", "System.setPassword() — AppExchange violation",
                                     "Hardwired crypto keys", "CRUD violations (PMD)", "Database.executeAnonymous()", "Open redirect risk"]),
    (HIGH_ORANGE, "High (20)", ["Classes without sharing declaration", "Weak crypto (MD5/SHA-1)", "@IsTest(SeeAllData=true)",
                                 "No assert statements in tests", "Missing PMD design rules", "UserInfo.getSessionId()"]),
    (MEDIUM_YELLOW, "Medium (13)", ["Outdated API versions (<v55)", "@future annotation usage", "System.debug statements", "Excessive parameters"]),
    (LOW_GREEN, "Low (6)", ["Test data without @TestSetup", "debug() without LoggingLevel", "Other PMD hygiene rules"]),
]
y = 1.82
for color, label, items in cq:
    add_rect(slide, 0.3, y, 0.18, 0.25, fill=color)
    add_text(slide, label, 0.55, y, 3.7, 0.28, size=11, bold=True, color=DARK_TEXT)
    y += 0.32
    for item in items:
        add_text(slide, f"  • {item}", 0.35, y, 3.95, 0.3, size=9.5, color=DARK_TEXT)
        y += 0.3
    y += 0.12

# Test Coverage column
add_rect(slide, 4.6, 1.28, 4.0, 5.85, fill=LIGHT_GREY)
add_rect(slide, 4.6, 1.28, 4.0, 0.42, fill=SF_BLUE)
add_text(slide, "Test Coverage — 7 Checks", 4.72, 1.3, 3.8, 0.38,
         size=13, bold=True, color=WHITE)

tc = [
    (CRITICAL_RED, "Classes / triggers with zero coverage", "Untested components block every deployment. Abstract classes with 0 executable lines excluded."),
    (HIGH_ORANGE, "Below 75% coverage (org-owned only)", "Managed package classes excluded. Identifies true coverage risk."),
    (HIGH_ORANGE, "Low test class ratio (< 30%)", "Too few test classes relative to production code volume."),
    (MEDIUM_YELLOW, "Triggers without dedicated test class", "By naming convention — triggers tested indirectly are higher risk."),
    (MEDIUM_YELLOW, "Test classes with no System.runAs()", "Multi-user and sharing scenarios never validated. PMD rule."),
    (LOW_GREEN, "Assert statements missing message", "Failures output 'Assertion failed' — no context. PMD rule."),
    (LOW_GREEN, "Deprecated testMethod keyword", "Replace with @isTest annotation. PMD rule."),
]
y = 1.82
for color, title, desc in tc:
    add_rect(slide, 4.7, y, 0.18, 0.25, fill=color)
    add_text(slide, title, 4.95, y, 3.5, 0.28, size=11, bold=True, color=DARK_TEXT)
    y += 0.3
    add_text(slide, desc, 4.75, y, 3.72, 0.48, size=9.5, color=MID_GREY, italic=True)
    y += 0.55

# Flow Quality column
add_rect(slide, 8.8, 1.28, 4.3, 5.85, fill=LIGHT_GREY)
add_rect(slide, 8.8, 1.28, 4.3, 0.42, fill=SF_DARK)
add_text(slide, "Flow Quality — 6 Checks", 8.92, 1.3, 4.1, 0.38,
         size=13, bold=True, color=WHITE)

fq = [
    (HIGH_ORANGE, "Active Process Builder flows", "ProcessType = Workflow — legacy automation on Salesforce's deprecation path. Migrate to record-triggered flows."),
    (HIGH_ORANGE, "DML operations inside loops", "Flows with both Loop and record operation elements risk hitting governor limits on large data sets."),
    (HIGH_ORANGE, "System Context Without Sharing", "Bypasses record-level security entirely. Privilege escalation risk if user data influences filters."),
    (LOW_GREEN, "System Context With Sharing", "Elevated privileges — confirm this is intentional. Switch to User context where possible."),
    (LOW_GREEN, "Obsolete flow versions (> 50)", "Deactivated versions accumulate after each activation. Slows Setup and clutters version history."),
    (LOW_GREEN, "Active flows with no description", "Undocumented flows are difficult to audit and maintain over time."),
]
y = 1.82
for color, title, desc in fq:
    add_rect(slide, 8.9, y, 0.18, 0.25, fill=color)
    add_text(slide, title, 9.15, y, 3.8, 0.28, size=11, bold=True, color=DARK_TEXT)
    y += 0.3
    add_text(slide, desc, 8.95, y, 4.02, 0.52, size=9.5, color=MID_GREY, italic=True)
    y += 0.6


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Accuracy & How It Works
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Accuracy-First Design", "Built to present to customers — every check validated against Salesforce platform behaviour")
footer(slide, 8)

accuracy_items = [
    ("Managed packages excluded", "Every query filters NamespacePrefix = null or cross-references org-owned IDs. Managed package Flows, Apex classes, queues, duplicate rules, record types, and remote sites never inflate findings."),
    ("Canary queries prevent false positives", "Features that aren't licensed (Knowledge, OmniStudio) are detected via explicit canary queries. Findings only fire when the feature is confirmed enabled — not when a query silently returns empty."),
    ("Silent failures surfaced", "Known Salesforce API limitations (AppDefinition profile-filtering, ServiceChannel permission requirements) are handled with corroborating signal checks rather than blind count thresholds."),
    ("Edition-aware limits", "Org limit thresholds (custom objects, API calls, storage) use the actual values from the Limits REST API — not hardcoded guesses. Enterprise = 2,000 objects; Developer = 400; etc."),
    ("Standard vs Tooling API routing", "FlowElement, WorkflowRule, AppDefinition, and ApexClass body scans use the Tooling API. Standard SOQL is used where appropriate. Wrong routing causes silent failures — every query is validated."),
    ("Inactive records excluded", "Active/IsActive/Status filters applied consistently — workflow rules, escalation rules, assignment rules, email service functions, email templates, and Apex classes all scoped to active only."),
]

cw = 6.2
for i, (title, desc) in enumerate(accuracy_items):
    col = i % 2
    row = i // 2
    l = 0.25 + col * (cw + 0.42)
    t = 1.35 + row * 1.9
    add_rect(slide, l, t, cw, 1.78, fill=LIGHT_GREY, line=SF_BLUE, line_w=0.4)
    add_rect(slide, l, t, 0.18, 1.78, fill=SF_BLUE)
    add_text(slide, title, l + 0.3, t + 0.08, cw - 0.4, 0.36,
             size=13, bold=True, color=SF_DARK)
    add_text(slide, desc, l + 0.3, t + 0.48, cw - 0.42, 1.22,
             size=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Sample Output / Scoring
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Assessment Output", "Scored, prioritised, and export-ready in under 60 seconds")
footer(slide, 9)

# Score visualisation
add_rect(slide, 0.25, 1.28, 3.5, 5.85, fill=SF_DARK)
add_text(slide, "Overall Score", 0.35, 1.38, 3.3, 0.4,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "62", 0.35, 1.9, 3.3, 1.2,
         size=72, bold=True, color=SF_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "/ 100", 0.35, 3.05, 3.3, 0.5,
         size=20, color=SF_LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Example org — illustrative", 0.35, 3.55, 3.3, 0.35,
         size=10, color=MID_GREY, italic=True, align=PP_ALIGN.CENTER)

sample_scores = [
    ("Sharing & Security", 48, CRITICAL_RED),
    ("Flow Quality", 55, CRITICAL_RED),
    ("Code Quality", 61, HIGH_ORANGE),
    ("Service Cloud", 67, HIGH_ORANGE),
    ("Test Coverage", 72, MEDIUM_YELLOW),
    ("Configuration", 79, MEDIUM_YELLOW),
    ("Integrations", 85, LOW_GREEN),
    ("Org Limits", 94, LOW_GREEN),
]

for i, (cat, score, color) in enumerate(sample_scores):
    y = 4.1 + i * 0.34
    bar_w = (score / 100) * 2.8
    add_rect(slide, 0.35, y, 2.8, 0.24, fill=RGBColor(0x30, 0x40, 0x70))
    add_rect(slide, 0.35, y, bar_w, 0.24, fill=color)
    add_text(slide, cat, 0.35, y, 2.4, 0.24, size=9, color=WHITE)
    add_text(slide, str(score), 3.1, y, 0.5, 0.24, size=9, bold=True, color=WHITE)

# Findings panel
add_rect(slide, 3.95, 1.28, 5.5, 5.85, fill=LIGHT_GREY)
add_rect(slide, 3.95, 1.28, 5.5, 0.38, fill=SF_BLUE)
add_text(slide, "Sample Findings (prioritised by severity)", 4.08, 1.3, 5.3, 0.34,
         size=12, bold=True, color=WHITE)

findings = [
    (CRITICAL_RED, "C", "213 Objects with Internal Public Read/Write OWD"),
    (CRITICAL_RED, "C", "47 Non-Admin Users with Modify All Data"),
    (CRITICAL_RED, "C", "12 Classes with CRUD violations (no FLS check)"),
    (HIGH_ORANGE, "H", "Active Process Builder Flows — legacy automation"),
    (HIGH_ORANGE, "H", "SOAP login() still in use — disabled Spring '26"),
    (HIGH_ORANGE, "H", "No Escalation Rules configured despite 130 queues"),
    (HIGH_ORANGE, "H", "Test classes not using System.runAs()"),
    (MEDIUM_YELLOW, "M", "873 → 241 Classes Below 75% Coverage (scoped)"),
    (MEDIUM_YELLOW, "M", "Stale published Knowledge articles (12+ months)"),
    (MEDIUM_YELLOW, "M", "Email templates include inactive/retired templates"),
    (LOW_GREEN, "L", "130 Queues — review routing consolidation"),
    (LOW_GREEN, "L", "Active flows with no description"),
]

for i, (color, sev, title) in enumerate(findings):
    y = 1.78 + i * 0.41
    add_rect(slide, 4.05, y, 0.28, 0.3, fill=color)
    add_text(slide, sev, 4.05, y + 0.01, 0.28, 0.28,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 4.4, y + 0.03, 4.9, 0.3, size=10, color=DARK_TEXT)

# Export options
add_rect(slide, 9.65, 1.28, 3.45, 5.85, fill=SF_DARK)
add_text(slide, "Export Options", 9.78, 1.35, 3.2, 0.38,
         size=14, bold=True, color=WHITE)

exports = [
    ("PDF Report", "Full category scores, all findings, recommendations, and affected record lists per issue."),
    ("Excel Workbook", "One tab per category + Summary tab. Filter and sort by severity."),
    ("CSV Flat File", "One row per affected record. Import directly into a project tracker or Jira."),
    ("Remediation Roadmap", "Full-screen phased view — Critical → High → Medium → Low. Print to PDF for stakeholders."),
]

y = 1.85
for title, desc in exports:
    add_rect(slide, 9.75, y, 3.2, 0.3, fill=SF_BLUE)
    add_text(slide, title, 9.82, y + 0.02, 3.1, 0.26,
             size=11, bold=True, color=WHITE)
    y += 0.36
    add_text(slide, desc, 9.82, y, 3.1, 0.62, size=10, color=SF_LIGHT_BLUE)
    y += 0.8


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — How to Use It (Setup)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Getting Started — 15-Minute Setup Per Org", "One Connected App registration. Then connect any time.")
footer(slide, 10)

# Steps
steps_detail = [
    ("Setup → App Manager", "New Connected App → enable OAuth → add callback URL:\nhttps://sf-tech-debt-assessor.onrender.com/auth/callback\nScopes: api + refresh_token  |  Disable PKCE  |  Save → wait 10 min"),
    ("Copy credentials", "App Manager → View → Manage Consumer Details\nCopy Consumer Key (Client ID) and Consumer Secret"),
    ("Connect", "Open sf-tech-debt-assessor.onrender.com\nEnter org URL, Client ID, Client Secret → Connect\nLog in with Salesforce → Allow"),
    ("Assessment runs automatically", "~45 seconds for a full org scan\nAll 341 checks run in parallel across 23 categories\nNo action required — just wait for results"),
]

for i, (title, detail) in enumerate(steps_detail):
    l = 0.25 + i * 3.25
    add_rect(slide, l, 1.35, 3.0, 0.52, fill=SF_BLUE)
    add_text(slide, f"Step {i+1}", l + 0.1, 1.38, 0.7, 0.26,
             size=11, color=SF_DARK, bold=True)
    add_text(slide, title, l + 0.1, 1.62, 2.8, 0.22,
             size=10, bold=True, color=WHITE)
    add_rect(slide, l, 1.87, 3.0, 3.5, fill=LIGHT_GREY)
    add_text(slide, detail, l + 0.12, 1.95, 2.78, 3.35,
             size=11, color=DARK_TEXT)

add_rect(slide, 0.25, 5.55, 12.8, 1.25, fill=SF_DARK)
add_text(slide, "Permissions required on the authenticating user:", 0.4, 5.6, 5, 0.35,
         size=12, bold=True, color=WHITE)
perms = ["API Enabled", "View Setup and Configuration", "Modify Metadata Through Metadata API Functions",
         "System Administrator profile provides all of the above"]
for i, p in enumerate(perms):
    add_text(slide, f"✓  {p}", 0.4 + (i % 2) * 6.4, 5.98 + (i // 2) * 0.35, 6.0, 0.32,
             size=11, color=SF_LIGHT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Why This vs Manual Assessment
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Why Use This Tool", "Faster, deeper, and more accurate than manual org review")
footer(slide, 11)

# Comparison table
headers = ["", "Manual Assessment", "SF Tech Debt Assessor"]
rows = [
    ("Time to complete", "2–3 days", "< 60 seconds"),
    ("Categories covered", "Varies (5–10 typical)", "23 categories, every run"),
    ("Check depth", "~50–80 checks", "341 checks"),
    ("Managed packages", "Often included (inflated)", "Excluded — namespace-filtered"),
    ("Inactive records", "Often included", "Filtered — active only"),
    ("Export formats", "Manual document", "PDF, Excel, CSV, Roadmap"),
    ("Repeatability", "Single point in time", "Re-run any time, track over time"),
    ("Silent failures", "Unknown", "Canary queries detect unavailable APIs"),
    ("Cost", "Consultant hours", "Free — self-service"),
]

col_widths = [3.8, 4.2, 4.7]
col_starts = [0.25, 4.15, 8.45]
row_h = 0.48
header_y = 1.3

for ci, (hdr, w, l) in enumerate(zip(headers, col_widths, col_starts)):
    bg = SF_DARK if ci == 0 else (CRITICAL_RED if ci == 1 else SF_BLUE)
    add_rect(slide, l, header_y, w, 0.44, fill=bg)
    add_text(slide, hdr, l + 0.1, header_y + 0.05, w - 0.15, 0.34,
             size=13, bold=True, color=WHITE)

for ri, row_data in enumerate(rows):
    y = header_y + 0.44 + ri * row_h
    bg_row = LIGHT_GREY if ri % 2 == 0 else WHITE
    for ci, (val, w, l) in enumerate(zip(row_data, col_widths, col_starts)):
        bg = bg_row if ci > 0 else SF_DARK
        color = WHITE if ci == 0 else (MID_GREY if ci == 1 else SF_DARK)
        add_rect(slide, l, y, w, row_h - 0.04, fill=bg)
        add_text(slide, val, l + 0.1, y + 0.06, w - 0.15, row_h - 0.1,
                 size=11, color=color, bold=(ci == 2))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Closing / Call to Action
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=SF_DARK)
add_rect(slide, 0, 0, 13.33, 0.1, fill=SF_BLUE)
add_rect(slide, 0, 7.4, 13.33, 0.1, fill=SF_BLUE)

add_text(slide, "Ready to assess your org?", 0.5, 1.1, 12.3, 0.8,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "sf-tech-debt-assessor.onrender.com",
         0.5, 2.05, 12.3, 0.65,
         size=28, bold=True, color=SF_BLUE, align=PP_ALIGN.CENTER)

add_text(slide, "15-minute setup · Read-only · Free · No data stored",
         0.5, 2.9, 12.3, 0.45,
         size=15, color=SF_LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)

ctas = [
    ("341 Checks", "Across 23 categories — the most\ncomprehensive Salesforce org audit available."),
    ("Read-Only", "Zero writes to Salesforce.\nAll results live in your browser tab only."),
    ("Export Ready", "PDF, Excel, CSV, and Remediation\nRoadmap — ready to share with stakeholders."),
    ("Always Current", "Built against Spring '26 / Summer '26\nenforcement dates and release notes."),
]

for i, (title, desc) in enumerate(ctas):
    l = 0.5 + i * 3.1
    add_rect(slide, l, 3.65, 2.9, 2.1, fill=RGBColor(0x1E, 0x35, 0x7A))
    add_rect(slide, l, 3.65, 2.9, 0.4, fill=SF_BLUE)
    add_text(slide, title, l + 0.1, 3.67, 2.7, 0.36,
             size=14, bold=True, color=WHITE)
    add_text(slide, desc, l + 0.1, 4.12, 2.7, 1.55,
             size=12, color=SF_LIGHT_BLUE)

add_text(slide, "Steven Bilgram  |  Success Architect  |  sbilgram@lgtm.io",
         0.5, 6.1, 12.3, 0.4,
         size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "github.com/sbilgram-lgtm/sf-tech-debt-assessor",
         0.5, 6.5, 12.3, 0.4,
         size=12, color=SF_LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)


# ── Save ────────────────────────────────────────────────────────────────────
out = "docs/SF_Tech_Debt_Assessor_Executive_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
