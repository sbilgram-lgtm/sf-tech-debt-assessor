"""
Generates an executive PowerPoint presentation for SF Tech Debt Assessor.
Saves to ~/Desktop/SF_Tech_Debt_Assessor_Executive_Presentation_2026-06-10.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Brand colours ────────────────────────────────────────────────
NAVY      = RGBColor(0x03, 0x2D, 0x60)   # Salesforce dark navy
SF_BLUE   = RGBColor(0x00, 0x70, 0xD2)   # Salesforce blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
BLUE      = RGBColor(0x34, 0x98, 0xDB)
ORANGE    = RGBColor(0xD3, 0x54, 0x00)
RED       = RGBColor(0xC0, 0x39, 0x2B)
PURPLE    = RGBColor(0x8E, 0x44, 0xAD)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF7)
MID_GREY  = RGBColor(0x7F, 0x8C, 0x8D)
DARK_GREY = RGBColor(0x2C, 0x3E, 0x50)
AMBER     = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=1):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, lines, l, t, w, h, size=13, color=NAVY, bold_first=False, spacing=None):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if spacing:
            p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)

def header_band(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.2, fill=NAVY)
    add_text(slide, title, 0.4, 0.12, 10, 0.6,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 10, 0.45,
                 size=14, color=RGBColor(0xAB, 0xB2, 0xB9))

def footer(slide, num):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=NAVY)
    add_text(slide, "SF Tech Debt Assessor  |  Confidential", 0.3, 7.21, 8, 0.28,
             size=9, color=RGBColor(0xAB, 0xB2, 0xB9))
    add_text(slide, str(num), 12.8, 7.21, 0.4, 0.28,
             size=9, color=RGBColor(0xAB, 0xB2, 0xB9), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# Salesforce-gradient background
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(slide, 0, 5.8, 13.33, 1.7, fill=SF_BLUE)

add_text(slide, "SF Tech Debt Assessor", 0.6, 1.4, 12, 1.1,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Automated Salesforce Org Health Assessment", 0.6, 2.6, 12, 0.6,
         size=22, color=RGBColor(0xAB, 0xB2, 0xB9), align=PP_ALIGN.CENTER)
add_text(slide, "22 Categories  ·  311 Checks  ·  Read-Only OAuth", 0.6, 3.2, 12, 0.5,
         size=15, color=RGBColor(0xAB, 0xB2, 0xB9), align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "Steven Bilgram, Success Architect  |  2026",
         0.6, 6.05, 12, 0.5,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "The Problem", "Org health reviews are manual, inconsistent, and time-consuming")
footer(slide, 2)

pain_points = [
    ("⏱  4–8 Hours of Manual Work",
     "Each assessment requires navigating Setup, running ad-hoc SOQL queries,\n"
     "and writing a findings document by hand."),
    ("📋  No Standardised Methodology",
     "Coverage depends entirely on individual consultant experience.\n"
     "~60–70% of risk areas are reviewed in a typical manual assessment."),
    ("🔍  No Drill-Down to Root Cause",
     "Findings are high-level. Clients cannot see which specific records,\n"
     "rules, or users are causing the score reduction."),
    ("📂  No Structured Remediation Path",
     "Recommendations are delivered as narrative text — rarely assigned,\n"
     "rarely actioned. ~30–40% of findings result in a task within 90 days."),
]

for i, (title, body) in enumerate(pain_points):
    col = 0.35 if i % 2 == 0 else 6.85
    top = 1.45 if i < 2 else 4.0
    add_rect(slide, col, top, 6.0, 2.3, fill=WHITE, line=RED)
    add_text(slide, title, col + 0.2, top + 0.15, 5.6, 0.45,
             size=13, bold=True, color=RED)
    add_text(slide, body, col + 0.2, top + 0.62, 5.6, 1.5,
             size=12, color=NAVY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "The Solution", "SF Tech Debt Assessor — from zero to scored report in under 5 minutes")
footer(slide, 3)

add_text(slide,
         "A web app that connects directly to any Salesforce org via OAuth and automatically "
         "runs a scored technical debt assessment across 22 categories and 311 checks — producing a "
         "stakeholder-ready report, drill-down record detail, and a phased remediation roadmap.",
         0.4, 1.3, 12.5, 0.9, size=14, color=NAVY)

steps = [
    ("1", "Connect",  "Authenticate via OAuth\nusing your org credentials"),
    ("2", "Assess",   "22-category scan across\n311 checks runs automatically"),
    ("3", "Review",   "Scored report with\ndrill-down to affected records"),
    ("4", "Export",   "PDF · Excel · CSV ·\nRemediation Roadmap"),
]

for i, (num, title, body) in enumerate(steps):
    x = 0.35 + i * 3.17
    add_rect(slide, x, 2.45, 2.85, 2.8, fill=NAVY)
    add_text(slide, num, x + 0.2, 2.55, 0.5, 0.55,
             size=28, bold=True, color=SF_BLUE)
    add_text(slide, title, x + 0.2, 3.15, 2.5, 0.45,
             size=15, bold=True, color=WHITE)
    add_text(slide, body, x + 0.2, 3.65, 2.5, 1.4,
             size=12, color=RGBColor(0xAB, 0xB2, 0xB9))
    if i < 3:
        add_text(slide, "→", x + 2.88, 3.4, 0.3, 0.5,
                 size=20, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

add_text(slide, "Deployed at: sf-tech-debt-assessor.onrender.com",
         0.4, 5.55, 12.5, 0.4,
         size=12, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — 21 Assessment Categories (with check counts)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "What It Assesses",
            "22 categories · 311 checks covering the full Salesforce technical stack")
footer(slide, 4)

categories = [
    ("Configuration",             13),
    ("Code Quality",              29),
    ("Data Model",                 4),
    ("Service Cloud",             69),
    ("Sharing & Security",        24),
    ("Integrations",               9),
    ("Test Coverage",              4),
    ("Org Limits",                 5),
    ("Duplicate & Matching Rules", 4),
    ("Reports & Dashboards",       3),
    ("Email Templates",            3),
    ("Platform Events & CDC",      3),
    ("Managed Packages",           3),
    ("Custom Metadata & Settings", 3),
    ("Record Types & Page Layouts",4),
    ("Einstein & AI Usage",        9),
    ("Experience Cloud",          12),
    ("Connected App Security",    12),
    ("LWC & Components",          39),
    ("OmniStudio",                26),
    ("Performance",               21),
    ("Notes & Attachments",       12),
]

cols = 3
col_w = 4.1
row_h = 0.53
for i, (cat, count) in enumerate(categories):
    col = i % cols
    row = i // cols
    x = 0.35 + col * col_w
    y = 1.38 + row * row_h
    add_rect(slide, x, y, col_w - 0.15, row_h - 0.07, fill=WHITE, line=SF_BLUE)
    add_text(slide, f"✓  {cat}", x + 0.15, y + 0.05, col_w - 1.2, 0.36,
             size=11, color=NAVY)
    add_text(slide, f"{count} checks", x + col_w - 1.28, y + 0.07, 1.0, 0.3,
             size=9.5, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — UX & Experience Highlights (NEW)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "User Experience",
            "Designed for consultants, admins, and customers — minimal setup, instant insight")
footer(slide, 5)

ux_features = [
    (SF_BLUE,  "Branded Split-Screen Landing Page",
     "Salesforce blue gradient hero panel with headline, stats bar, and 21-category grid "
     "displayed alongside the OAuth credential form — value proposition visible at the moment of login."),
    (PURPLE,   "Category Detail Modals",
     "Click any category card on the landing page to instantly see all checks for that "
     "category, each with a color-coded severity badge (Critical · High · Medium · Low)."),
    (GREEN,    "In-App Setup Guide",
     "'Need Instructions?' opens a full step-by-step Connected App / External Client App "
     "setup walkthrough — including troubleshooting — without leaving the app."),
    (SF_BLUE,  "Auto-Run on Login",
     "The assessment starts immediately after OAuth completion — no intermediate screen, "
     "no extra button click. Progress is shown category by category as the scan runs."),
    (ORANGE,   "Drill-Down to Affected Records",
     "Every finding expands to show the exact classes, rules, users, fields, or pages "
     "causing the score deduction — not just high-level categories."),
    (PURPLE,   "4-Phase Remediation Roadmap",
     "Full-screen phased action plan ordered by severity then category. "
     "Every item includes title, description, recommendation, and record count. "
     "Print or save as PDF in one click."),
]

for i, (color, title, body) in enumerate(ux_features):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.45
    y = 1.38 + row * 1.85
    add_rect(slide, x, y, 6.1, 1.68, fill=WHITE, line=color)
    add_rect(slide, x, y, 0.18, 1.68, fill=color)
    add_text(slide, title, x + 0.35, y + 0.1, 5.6, 0.38,
             size=12.5, bold=True, color=color)
    add_text(slide, body, x + 0.35, y + 0.52, 5.6, 1.05,
             size=11, color=NAVY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Value Proposition
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "Value Proposition", "What changes for every persona who uses the tool")
footer(slide, 6)

personas = [
    (SF_BLUE, "Consultant / SA",
     "Delivers a scored assessment in the kickoff meeting.\nEliminates 4–8 hrs of manual work per engagement."),
    (GREEN,   "Salesforce Admin",
     "Gains visibility into hidden risk.\nSelf-sufficient — no consultant required."),
    (PURPLE,  "IT Manager / Director",
     "Receives evidence-backed, export-ready reports.\nFaster decisions, clearer remediation investment."),
    (ORANGE,  "Practice Lead",
     "Standardises delivery across the team.\nEnables a repeatable assessment service offering."),
    (NAVY,    "Customer Success Mgr",
     "Identifies at-risk customers before issues escalate.\nElevates QBRs with objective org health scores."),
    (RED,     "End Client / Stakeholder",
     "Understands risk without deep Salesforce expertise.\nReceives a phased action plan on day one."),
]

for i, (color, persona, desc) in enumerate(personas):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.45
    y = 1.4 + row * 1.85
    add_rect(slide, x, y, 6.1, 1.65, fill=WHITE, line=color)
    add_rect(slide, x, y, 0.18, 1.65, fill=color)
    add_text(slide, persona, x + 0.35, y + 0.12, 5.6, 0.38,
             size=13, bold=True, color=color)
    add_text(slide, desc, x + 0.35, y + 0.52, 5.6, 1.0,
             size=11.5, color=NAVY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Quantified Impact
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "Quantified Impact", "Measured against the manual assessment baseline")
footer(slide, 7)

metrics = [
    (GREEN,   "6–8 hrs",    "Saved per\nassessment"),
    (SF_BLUE, "60–70%",     "Reduction in\ntime to first action"),
    (PURPLE,  "3–5×",       "Increase in\nassessment frequency"),
    (ORANGE,  "$900–$2K",   "Cost savings\nper engagement"),
    (RED,     "20–30%",     "Reduction in reactive\nsupport cases / quarter"),
    (NAVY,    "100%",       "Category coverage\nevery time"),
]

for i, (color, value, label) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.3
    y = 1.45 + row * 2.65
    add_rect(slide, x, y, 3.95, 2.25, fill=WHITE, line=color)
    add_text(slide, value, x + 0.2, y + 0.2, 3.55, 0.85,
             size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.2, y + 1.1, 3.55, 0.9,
             size=12.5, color=NAVY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Outputs
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header_band(slide, "Key Outputs", "Every assessment produces four ready-to-use artifacts")
footer(slide, 8)

outputs = [
    (SF_BLUE, "Export PDF",
     "Full scored report with category breakdowns, findings, recommendations, and "
     "affected records per issue. Includes Org Name, Org ID, Type, Instance, and URL. "
     "Ideal for customer presentations and stakeholder reviews."),
    (GREEN,   "Export Excel",
     "One tab per category plus a Summary tab. Each tab lists all findings with "
     "affected records inline. Includes org metadata on every tab. "
     "Opens directly in Excel for further analysis or import."),
    (GREEN,   "Export CSV",
     "Flat file with one row per affected record across all categories. "
     "Includes org identifier columns on every row. "
     "Ideal for importing into a project tracker, Jira, or remediation backlog."),
    (PURPLE,  "Remediation Roadmap",
     "Full-screen 4-phase action plan ordered by severity (Critical → High → Medium → Low), "
     "then by category. Each item includes title, description, recommended action, "
     "and affected record count. Print or save as PDF in one click."),
]

for i, (color, title, body) in enumerate(outputs):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.45
    y = 1.4 + row * 2.6
    add_rect(slide, x, y, 6.1, 2.35, fill=WHITE, line=color)
    add_rect(slide, x, y, 6.1, 0.48, fill=color)
    add_text(slide, title, x + 0.2, y + 0.07, 5.7, 0.38,
             size=14, bold=True, color=WHITE)
    add_text(slide, body, x + 0.2, y + 0.65, 5.7, 1.55,
             size=11.5, color=NAVY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Call to Action
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(slide, 0, 5.8, 13.33, 1.7, fill=SF_BLUE)

add_text(slide, "Ready to See It in Action?", 0.6, 1.1, 12, 0.9,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide,
         "SF Tech Debt Assessor is live and available today.\n"
         "Connect any Salesforce org and run a full 22-category, 311-check assessment in under 5 minutes.",
         0.6, 2.1, 12, 0.9,
         size=16, color=RGBColor(0xAB, 0xB2, 0xB9), align=PP_ALIGN.CENTER)

# Stat bar
stats = [("22", "Categories"), ("311", "Checks"), ("100%", "Read-Only")]
for i, (val, lbl) in enumerate(stats):
    x = 1.5 + i * 3.5
    add_rect(slide, x, 3.2, 3.0, 1.3, fill=RGBColor(0x07, 0x50, 0x9A))
    add_text(slide, val,  x, 3.3, 3.0, 0.65,
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, lbl,  x, 3.95, 3.0, 0.4,
             size=11, color=RGBColor(0xAB, 0xB2, 0xB9), align=PP_ALIGN.CENTER)

add_rect(slide, 3.67, 4.85, 6.0, 0.72, fill=SF_BLUE)
add_text(slide, "sf-tech-debt-assessor.onrender.com",
         3.67, 4.88, 6.0, 0.65,
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Built by Steven Bilgram, Success Architect",
         0.6, 6.05, 12, 0.4,
         size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────
out = os.path.expanduser(
    "~/Desktop/SF_Tech_Debt_Assessor_Executive_Presentation_2026-06-10.pptx"
)
prs.save(out)
print(f"Saved: {out}")
